from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # Create presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Define color scheme
    PRIMARY_COLOR = RGBColor(37, 99, 235)  # Blue
    SECONDARY_COLOR = RGBColor(71, 85, 105)  # Dark gray
    ACCENT_COLOR = RGBColor(16, 185, 129)  # Green

    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "ChatBot with PDF"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "AI-Powered Document Intelligence System"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
    subtitle_para.alignment = PP_ALIGN.CENTER

    # Author info
    author_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.5))
    author_frame = author_box.text_frame
    author_frame.text = "ISE492 Project Presentation"
    author_para = author_frame.paragraphs[0]
    author_para.font.size = Pt(18)
    author_para.font.color.rgb = RGBColor(255, 255, 255)
    author_para.alignment = PP_ALIGN.CENTER

    # Slide 2: Project Overview
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Project Overview"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.5), Inches(4.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    content_items = [
        "An intelligent chatbot that allows users to have conversations with their documents",
        "Supports multiple formats: PDF, DOCX, PPTX, and images",
        "Uses AI to understand document content and provide accurate answers",
        "Features include multi-document chat, folders, and precise page citations",
        "Built with modern full-stack technologies and RAG (Retrieval-Augmented Generation)"
    ]

    for i, item in enumerate(content_items):
        p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
        p.text = f"• {item}"
        p.font.size = Pt(20)
        p.font.color.rgb = SECONDARY_COLOR
        p.space_before = Pt(12)
        p.level = 0

    # Slide 3: System Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "System Architecture"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR

    # Architecture boxes
    # Frontend
    frontend_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(4), Inches(1.8))
    frontend_frame = frontend_box.text_frame
    frontend_frame.text = "Frontend\n\n• React + Vite\n• Tailwind CSS\n• PDF.js Viewer\n• Responsive UI"
    for p in frontend_frame.paragraphs:
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)

    frontend_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.8), Inches(1.8), Inches(4), Inches(1.8)
    )
    fill = frontend_shape.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    frontend_shape.text_frame.text = "Frontend\n\n• React + Vite\n• Tailwind CSS\n• PDF.js Viewer\n• Responsive UI"
    for p in frontend_shape.text_frame.paragraphs:
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.bold = True if p == frontend_shape.text_frame.paragraphs[0] else False

    # Backend
    backend_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(5.2), Inches(1.8), Inches(4), Inches(1.8)
    )
    fill = backend_shape.fill
    fill.solid()
    fill.fore_color.rgb = ACCENT_COLOR
    backend_shape.text_frame.text = "Backend\n\n• Node.js + Express\n• MongoDB\n• OpenAI API\n• JWT Auth"
    for p in backend_shape.text_frame.paragraphs:
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.bold = True if p == backend_shape.text_frame.paragraphs[0] else False

    # Database
    db_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.8), Inches(4.2), Inches(4), Inches(1.5)
    )
    fill = db_shape.fill
    fill.solid()
    fill.fore_color.rgb = SECONDARY_COLOR
    db_shape.text_frame.text = "Database\n\n• User Documents\n• Conversations\n• Vector Embeddings"
    for p in db_shape.text_frame.paragraphs:
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.bold = True if p == db_shape.text_frame.paragraphs[0] else False

    # AI Services
    ai_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(5.2), Inches(4.2), Inches(4), Inches(1.5)
    )
    fill = ai_shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(168, 85, 247)  # Purple
    ai_shape.text_frame.text = "AI Services\n\n• GPT-3.5 Turbo\n• Text Embeddings\n• GPT-4 Vision"
    for p in ai_shape.text_frame.paragraphs:
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.bold = True if p == ai_shape.text_frame.paragraphs[0] else False

    # Slide 4: Document Upload Workflow
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "How It Works: Document Upload"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR

    # Workflow steps
    steps = [
        ("1. Upload", "User uploads PDF, DOCX, PPTX, or image file"),
        ("2. Processing", "Extract text using pdf-parse, Mammoth, or Tesseract OCR"),
        ("3. Chunking", "Split text into 800-token chunks with 100-token overlap"),
        ("4. Embedding", "Generate vector embeddings using OpenAI API"),
        ("5. Storage", "Save document metadata and embeddings to MongoDB"),
        ("6. Ready", "Document is now ready for intelligent Q&A")
    ]

    y_position = 1.8
    for step_num, (title, desc) in enumerate(steps):
        # Step box
        box_shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(1), Inches(y_position), Inches(8), Inches(0.7)
        )
        fill = box_shape.fill
        fill.solid()
        fill.fore_color.rgb = PRIMARY_COLOR if step_num % 2 == 0 else ACCENT_COLOR

        text_frame = box_shape.text_frame
        text_frame.text = f"{title}: {desc}"
        p = text_frame.paragraphs[0]
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.bold = True

        y_position += 0.85

    # Slide 5: Chat Functionality - Step by Step
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "How It Works: Chat Processing"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR

    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.5), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    chat_steps = [
        "1. User types a question about the document",
        "2. System classifies question (document-related or general knowledge)",
        "3. Detects page references (e.g., 'page 5', 'slide 10')",
        "4. Generates embedding vector for the question",
        "5. Searches document embeddings using cosine similarity",
        "6. Retrieves top 15 most relevant text chunks",
        "7. Builds context with chunk content + page numbers",
        "8. Sends context + question to GPT-3.5 Turbo",
        "9. AI generates answer with mandatory page citations",
        "10. Returns response with clickable citations"
    ]

    for i, step in enumerate(chat_steps):
        p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
        p.text = step
        p.font.size = Pt(18)
        p.font.color.rgb = SECONDARY_COLOR
        p.space_before = Pt(8)

    # Slide 6: RAG System Explained
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "RAG System (Retrieval-Augmented Generation)"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(38)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR

    # Explanation
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.5), Inches(2))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    text_frame.text = "RAG combines information retrieval with AI generation to provide accurate, context-aware answers from your documents."
    p = text_frame.paragraphs[0]
    p.font.size = Pt(20)
    p.font.color.rgb = SECONDARY_COLOR
    p.font.italic = True

    # RAG Components
    components_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(8.5), Inches(3))
    text_frame = components_box.text_frame
    text_frame.word_wrap = True

    rag_components = [
        "Retrieval: Semantic search finds relevant document sections",
        "  • Uses vector embeddings (1536 dimensions)",
        "  • Cosine similarity matching (threshold: 0.7)",
        "  • Returns top 15 chunks ranked by relevance",
        "",
        "Augmentation: Retrieved content enhances AI's knowledge",
        "  • Adds document context to the AI prompt",
        "  • Includes page numbers and metadata",
        "",
        "Generation: AI creates accurate, cited responses",
        "  • GPT-3.5 Turbo with temperature=0 (deterministic)",
        "  • Mandatory citation format enforced"
    ]

    for i, item in enumerate(rag_components):
        p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
        p.text = item
        if item.startswith("  •"):
            p.font.size = Pt(16)
            p.level = 1
        elif item == "":
            continue
        else:
            p.font.size = Pt(18)
            p.font.bold = True
            p.level = 0
        p.font.color.rgb = SECONDARY_COLOR

    # Slide 7: Multi-Document Chat Feature
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Multi-Document Chat"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR

    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.5), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    multi_doc_features = [
        "Folder Organization",
        "  • Group related documents into folders",
        "  • Drag-and-drop document management",
        "",
        "Parallel Search",
        "  • Searches all documents in folder simultaneously",
        "  • Retrieves top 10 chunks from each document",
        "  • Selects best 15 chunks across all documents",
        "",
        "Source Attribution",
        "  • Citations include document name + page number",
        "  • Example: [research_paper.pdf - Page 12]",
        "  • Click citation to view exact source",
        "",
        "Smart Document Mapping",
        "  • Ensures citations reference correct document",
        "  • Handles MongoDB query ordering correctly"
    ]

    for i, item in enumerate(multi_doc_features):
        p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
        p.text = item
        if item.startswith("  •") or item.startswith("  Example:"):
            p.font.size = Pt(16)
            p.level = 1
        elif item == "":
            continue
        else:
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = ACCENT_COLOR
            p.level = 0
        p.font.color.rgb = SECONDARY_COLOR if not item or item.startswith(" ") else ACCENT_COLOR

    # Slide 8: Citation System
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Intelligent Citation System"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR

    # Citation formats
    format_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.5), Inches(2.5))
    text_frame = format_box.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = "Three Citation Formats:"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_COLOR

    citation_formats = [
        "Single Document: [Page 5] or [Slide 10]",
        "Multi-Document: [filename.pdf - Page 5]",
        "Natural Language: 'as mentioned on page 31...'"
    ]

    for item in citation_formats:
        p = text_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = SECONDARY_COLOR
        p.space_before = Pt(10)

    # Features
    features_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.3), Inches(8.5), Inches(2.5))
    text_frame = features_box.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = "Key Features:"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_COLOR

    citation_features = [
        "All citations are clickable - jump directly to source page",
        "AI is required to include citations in every response",
        "Page numbers extracted from document metadata",
        "Regex parsing identifies and formats citations"
    ]

    for item in citation_features:
        p = text_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = SECONDARY_COLOR
        p.space_before = Pt(10)

    # Slide 9: Key Technologies
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Technology Stack"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR

    # Two columns
    # Left column - Frontend
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(4), Inches(5))
    text_frame = left_box.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = "Frontend"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR

    frontend_tech = [
        "React 18",
        "Vite (Build Tool)",
        "Tailwind CSS",
        "PDF.js",
        "Mammoth.js (DOCX)",
        "Framer Motion",
        "React Router",
        "Axios"
    ]

    for tech in frontend_tech:
        p = text_frame.add_paragraph()
        p.text = f"• {tech}"
        p.font.size = Pt(16)
        p.font.color.rgb = SECONDARY_COLOR

    # Right column - Backend
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.6), Inches(4), Inches(5))
    text_frame = right_box.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = "Backend & AI"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR

    backend_tech = [
        "Node.js + Express",
        "MongoDB + Mongoose",
        "OpenAI GPT-3.5 Turbo",
        "OpenAI Embeddings",
        "GPT-4 Vision",
        "JWT Authentication",
        "pdf-parse",
        "Tesseract.js (OCR)"
    ]

    for tech in backend_tech:
        p = text_frame.add_paragraph()
        p.text = f"• {tech}"
        p.font.size = Pt(16)
        p.font.color.rgb = SECONDARY_COLOR

    # Slide 10: Output Example - Chat Interface
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Output: Chat Interface"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR

    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.5), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    features = [
        "Split-View Layout",
        "  • Document viewer on left, chat on right",
        "  • Resizable panels for flexible viewing",
        "",
        "Real-Time Responses",
        "  • Type question and get instant AI-generated answer",
        "  • Loading indicator during processing",
        "",
        "Message History",
        "  • Complete conversation history saved",
        "  • Context-aware responses using last 10 messages",
        "",
        "Interactive Features",
        "  • Click citations to navigate to source pages",
        "  • Markdown formatting support",
        "  • Copy/paste functionality"
    ]

    for i, item in enumerate(features):
        p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
        p.text = item
        if item.startswith("  •"):
            p.font.size = Pt(16)
            p.level = 1
        elif item == "":
            continue
        else:
            p.font.size = Pt(20)
            p.font.bold = True
            p.level = 0
        p.font.color.rgb = SECONDARY_COLOR

    # Slide 11: Output Example - Document Management
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Output: Document Management"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR

    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.5), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    doc_features = [
        "Document Library",
        "  • View all uploaded documents",
        "  • Shows file name, page count, upload date",
        "  • Search and filter capabilities",
        "",
        "Folder System",
        "  • Create folders to organize documents",
        "  • Drag-and-drop to move documents",
        "  • Rename, delete, and manage folders",
        "",
        "Document Actions",
        "  • Rename documents",
        "  • Delete documents",
        "  • Reset chat history",
        "  • Download original file",
        "",
        "Multi-Format Support",
        "  • PDF, DOCX, PPTX, and images",
        "  • Automatic format detection and rendering"
    ]

    for i, item in enumerate(doc_features):
        p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
        p.text = item
        if item.startswith("  •"):
            p.font.size = Pt(16)
            p.level = 1
        elif item == "":
            continue
        else:
            p.font.size = Pt(20)
            p.font.bold = True
            p.level = 0
        p.font.color.rgb = SECONDARY_COLOR

    # Slide 12: Special Features
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Advanced Features"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR

    # Feature boxes
    y_pos = 1.8
    advanced_features = [
        ("Image Extraction & Captioning", "Extracts images from PDFs and generates searchable captions using GPT-4 Vision"),
        ("Page-Specific Retrieval", "Detects 'page 5' or 'slide 10' in questions and retrieves only relevant pages"),
        ("Question Classification", "Determines if question is document-related or general knowledge"),
        ("Multi-Language Support", "English and Turkish language interface with auto-detection"),
        ("Secure Authentication", "JWT-based authentication with bcrypt password hashing")
    ]

    for title, desc in advanced_features:
        box_shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0.8), Inches(y_pos), Inches(8.5), Inches(0.85)
        )
        fill = box_shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(241, 245, 249)  # Light gray

        text_frame = box_shape.text_frame
        text_frame.word_wrap = True

        p = text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR

        p2 = text_frame.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(14)
        p2.font.color.rgb = SECONDARY_COLOR

        y_pos += 1

    # Slide 13: System Performance
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "System Performance & Accuracy"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR

    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.5), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    performance_points = [
        "Retrieval Accuracy",
        "  • Cosine similarity threshold: 0.7",
        "  • Top 15 most relevant chunks retrieved",
        "  • Average response time: 2-3 seconds",
        "",
        "Citation Accuracy",
        "  • 100% citation requirement enforced",
        "  • Page numbers verified against source documents",
        "  • Multi-document source attribution working correctly",
        "",
        "Scalability",
        "  • Vector embeddings cached in database",
        "  • Batch embedding generation (100 per API call)",
        "  • Async processing for document uploads",
        "",
        "User Experience",
        "  • Responsive design for all screen sizes",
        "  • Real-time loading indicators",
        "  • Smooth transitions and animations"
    ]

    for i, item in enumerate(performance_points):
        p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
        p.text = item
        if item.startswith("  •"):
            p.font.size = Pt(16)
            p.level = 1
        elif item == "":
            continue
        else:
            p.font.size = Pt(20)
            p.font.bold = True
            p.level = 0
        p.font.color.rgb = SECONDARY_COLOR

    # Slide 14: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(241, 245, 249)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Conclusion"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR
    title_para.alignment = PP_ALIGN.CENTER

    # Key achievements
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(3.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    conclusions = [
        "Successfully implemented a production-ready AI chatbot system",
        "",
        "Demonstrates full-stack development with modern technologies",
        "",
        "Advanced RAG system with semantic search and citations",
        "",
        "Multi-document chat with accurate source attribution",
        "",
        "Enterprise features: authentication, folders, document management"
    ]

    for i, item in enumerate(conclusions):
        p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
        if item == "":
            p.text = ""
            continue
        p.text = f"• {item}"
        p.font.size = Pt(20)
        p.font.color.rgb = SECONDARY_COLOR
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(8)

    # Thank you
    thanks_box = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(8), Inches(0.8))
    thanks_frame = thanks_box.text_frame
    thanks_frame.text = "Thank You!"
    thanks_para = thanks_frame.paragraphs[0]
    thanks_para.font.size = Pt(36)
    thanks_para.font.bold = True
    thanks_para.font.color.rgb = PRIMARY_COLOR
    thanks_para.alignment = PP_ALIGN.CENTER

    # Save presentation
    prs.save('ChatBot_with_PDF_Presentation.pptx')
    print("Presentation created successfully: ChatBot_with_PDF_Presentation.pptx")
    print(f"Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    create_presentation()
