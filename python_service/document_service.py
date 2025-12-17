"""
Document Processing Microservice
FastAPI service for extracting text from PDF, DOCX, PPTX, and images using OCR
"""

from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation
import pytesseract
from PIL import Image
import io
import os
import tempfile
from pathlib import Path
from loguru import logger
from typing import Dict, Any, List
import traceback
from table_extractor import extract_tables_from_pdf

# Configure logging
logger.add("document_service.log", rotation="10 MB", level="INFO")

app = FastAPI(
    title="Document Processing Service",
    description="Microservice for extracting text from various document formats",
    version="1.0.0"
)

# CORS Configuration - Allow Node.js backend to access this service
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3600", "http://localhost:5173"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Configure Tesseract path (adjust based on your installation)
# On Windows, Tesseract is usually installed at: C:\Program Files\Tesseract-OCR\tesseract.exe
# On Linux/Mac: /usr/bin/tesseract or /usr/local/bin/tesseract
if os.name == 'nt':  # Windows
    tesseract_paths = [
        r"C:\Program Files\Tesseract-OCR\tesseract.exe",
        r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
        r"C:\Tesseract-OCR\tesseract.exe"
    ]
    for path in tesseract_paths:
        if os.path.exists(path):
            pytesseract.pytesseract.tesseract_cmd = path
            logger.info(f"Tesseract found at: {path}")
            break
    else:
        logger.warning("Tesseract not found in standard paths. OCR may not work.")
else:
    # Linux/Mac - usually in PATH
    pytesseract.pytesseract.tesseract_cmd = 'tesseract'


@app.get("/")
async def root():
    """Root endpoint - service information"""
    return {
        "service": "Document Processing Microservice",
        "version": "1.0.0",
        "status": "running",
        "supported_formats": ["PDF", "DOCX", "PPTX", "Images (JPG, PNG, GIF, BMP, TIFF)"]
    }


@app.get("/health")
async def health_check():
    """Health check endpoint for monitoring"""
    return {
        "status": "healthy",
        "service": "document-processor",
        "tesseract_available": _check_tesseract()
    }


def _check_tesseract() -> bool:
    """Check if Tesseract is available"""
    try:
        pytesseract.get_tesseract_version()
        return True
    except Exception:
        return False


@app.post("/extract/pdf")
async def extract_pdf(file: UploadFile = File(...)) -> JSONResponse:
    """
    Extract text from PDF files using pdfplumber

    Returns:
        - text: Extracted text content
        - pages: Number of pages in the PDF
        - success: Processing status
        - filename: Original filename
    """
    logger.info(f"Processing PDF file: {file.filename}")

    try:
        # Read file content
        content = await file.read()

        # Process with pdfplumber
        with pdfplumber.open(io.BytesIO(content)) as pdf:
            # Extract text from all pages
            text_parts = []
            for page_num, page in enumerate(pdf.pages, 1):
                try:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(f"--- Page {page_num} ---\n{page_text}")
                except Exception as e:
                    logger.warning(f"Error extracting page {page_num}: {str(e)}")
                    continue

            full_text = "\n\n".join(text_parts)
            page_count = len(pdf.pages)

        # Check if any text was extracted
        if not full_text.strip():
            logger.warning(f"No text extracted from PDF: {file.filename}. May be scanned/image-based.")
            return JSONResponse(
                status_code=200,
                content={
                    "text": "",
                    "pages": page_count,
                    "success": True,
                    "warning": "No text extracted. This may be a scanned PDF. Try /extract/ocr endpoint.",
                    "filename": file.filename
                }
            )

        logger.info(f"Successfully extracted {len(full_text)} characters from {page_count} pages")

        return JSONResponse(
            status_code=200,
            content={
                "text": full_text,
                "pages": page_count,
                "success": True,
                "filename": file.filename,
                "char_count": len(full_text)
            }
        )

    except Exception as e:
        logger.error(f"Error processing PDF {file.filename}: {str(e)}\n{traceback.format_exc()}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to process PDF: {str(e)}"
        )


@app.post("/extract/docx")
async def extract_docx(file: UploadFile = File(...)) -> JSONResponse:
    """
    Extract text from DOCX files using python-docx

    Returns:
        - text: Extracted text content
        - paragraphs: Number of paragraphs
        - success: Processing status
        - filename: Original filename
    """
    logger.info(f"Processing DOCX file: {file.filename}")

    try:
        # Read file content
        content = await file.read()

        # Process with python-docx
        doc = DocxDocument(io.BytesIO(content))

        # Extract text from paragraphs
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)

        # Extract text from tables
        table_texts = []
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    table_texts.append(row_text)

        # Combine all text
        full_text = "\n\n".join(paragraphs)
        if table_texts:
            full_text += "\n\n--- Tables ---\n" + "\n".join(table_texts)

        if not full_text.strip():
            logger.warning(f"No text extracted from DOCX: {file.filename}")
            return JSONResponse(
                status_code=200,
                content={
                    "text": "",
                    "paragraphs": 0,
                    "success": True,
                    "warning": "No text content found in document",
                    "filename": file.filename
                }
            )

        logger.info(f"Successfully extracted {len(full_text)} characters from {len(paragraphs)} paragraphs")

        return JSONResponse(
            status_code=200,
            content={
                "text": full_text,
                "paragraphs": len(paragraphs),
                "tables": len(table_texts),
                "success": True,
                "filename": file.filename,
                "char_count": len(full_text)
            }
        )

    except Exception as e:
        logger.error(f"Error processing DOCX {file.filename}: {str(e)}\n{traceback.format_exc()}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to process DOCX: {str(e)}"
        )


@app.post("/extract/pptx")
async def extract_pptx(file: UploadFile = File(...)) -> JSONResponse:
    """
    Extract text from PPTX files using python-pptx

    Returns:
        - text: Extracted text content from all slides
        - slides: Number of slides
        - success: Processing status
        - filename: Original filename
    """
    logger.info(f"Processing PPTX file: {file.filename}")

    try:
        # Read file content
        content = await file.read()

        # Process with python-pptx
        prs = Presentation(io.BytesIO(content))

        # Extract text from all slides
        slide_texts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = []

            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())

                # Extract text from tables in slides
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells)
                        if row_text.strip():
                            slide_content.append(row_text)

            if slide_content:
                slide_text = f"--- Slide {slide_num} ---\n" + "\n".join(slide_content)
                slide_texts.append(slide_text)

        full_text = "\n\n".join(slide_texts)

        if not full_text.strip():
            logger.warning(f"No text extracted from PPTX: {file.filename}")
            return JSONResponse(
                status_code=200,
                content={
                    "text": "",
                    "slides": len(prs.slides),
                    "success": True,
                    "warning": "No text content found in presentation",
                    "filename": file.filename
                }
            )

        logger.info(f"Successfully extracted {len(full_text)} characters from {len(prs.slides)} slides")

        return JSONResponse(
            status_code=200,
            content={
                "text": full_text,
                "slides": len(prs.slides),
                "success": True,
                "filename": file.filename,
                "char_count": len(full_text)
            }
        )

    except Exception as e:
        logger.error(f"Error processing PPTX {file.filename}: {str(e)}\n{traceback.format_exc()}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to process PPTX: {str(e)}"
        )


@app.post("/extract-tables")
async def extract_tables(file: UploadFile = File(...)) -> JSONResponse:
    """
    TASK B: Extract tables from PDF and convert to Markdown format

    This endpoint extracts all tables from a PDF document and returns them
    as Markdown-formatted text, which is more suitable for embedding and RAG.

    Returns:
        - tables: List of extracted tables with page numbers and Markdown format
        - total_tables: Number of tables found
        - success: Processing status
        - filename: Original filename
    """
    logger.info(f"Extracting tables from PDF: {file.filename}")

    try:
        # Read file content
        content = await file.read()

        # Extract tables using our table_extractor module
        extracted_tables = extract_tables_from_pdf(content)

        if not extracted_tables:
            logger.info(f"No tables found in PDF: {file.filename}")
            return JSONResponse(
                status_code=200,
                content={
                    "tables": [],
                    "total_tables": 0,
                    "success": True,
                    "message": "No tables found in document",
                    "filename": file.filename
                }
            )

        logger.info(f"Successfully extracted {len(extracted_tables)} tables from {file.filename}")

        return JSONResponse(
            status_code=200,
            content={
                "tables": extracted_tables,
                "total_tables": len(extracted_tables),
                "success": True,
                "filename": file.filename
            }
        )

    except Exception as e:
        logger.error(f"Error extracting tables from {file.filename}: {str(e)}\n{traceback.format_exc()}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to extract tables: {str(e)}"
        )


@app.post("/extract/ocr")
async def extract_image_ocr(file: UploadFile = File(...)) -> JSONResponse:
    """
    Extract text from images using Tesseract OCR
    Supports: JPG, JPEG, PNG, GIF, BMP, TIFF

    Returns:
        - text: Extracted text via OCR
        - confidence: OCR confidence score (if available)
        - success: Processing status
        - filename: Original filename
    """
    logger.info(f"Processing image with OCR: {file.filename}")

    try:
        # Check if Tesseract is available
        if not _check_tesseract():
            raise HTTPException(
                status_code=503,
                detail="Tesseract OCR is not available. Please install Tesseract-OCR."
            )

        # Read file content
        content = await file.read()

        # Open image with PIL
        image = Image.open(io.BytesIO(content))

        # Convert to RGB if necessary (for RGBA or other formats)
        if image.mode not in ('RGB', 'L'):
            image = image.convert('RGB')

        # Enhance image for better OCR (optional preprocessing)
        # Resize if image is too large
        max_dimension = 3000
        if max(image.size) > max_dimension:
            ratio = max_dimension / max(image.size)
            new_size = tuple(int(dim * ratio) for dim in image.size)
            image = image.resize(new_size, Image.Resampling.LANCZOS)
            logger.info(f"Resized image to {new_size}")

        # Perform OCR
        logger.info("Running Tesseract OCR...")
        text = pytesseract.image_to_string(image, lang='eng')

        # Get OCR confidence data (optional)
        try:
            ocr_data = pytesseract.image_to_data(image, output_type=pytesseract.Output.DICT)
            confidences = [conf for conf in ocr_data['conf'] if conf != -1]
            avg_confidence = sum(confidences) / len(confidences) if confidences else 0
        except Exception:
            avg_confidence = None

        if not text.strip():
            logger.warning(f"No text extracted from image: {file.filename}")
            return JSONResponse(
                status_code=200,
                content={
                    "text": "",
                    "success": True,
                    "warning": "No text detected in image",
                    "filename": file.filename
                }
            )

        logger.info(f"Successfully extracted {len(text)} characters via OCR")

        result = {
            "text": text,
            "success": True,
            "filename": file.filename,
            "char_count": len(text),
            "image_size": image.size
        }

        if avg_confidence is not None:
            result["confidence"] = round(avg_confidence, 2)

        return JSONResponse(status_code=200, content=result)

    except Exception as e:
        logger.error(f"Error processing image {file.filename}: {str(e)}\n{traceback.format_exc()}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to process image: {str(e)}"
        )


@app.post("/convert/pptx-to-pdf")
async def convert_pptx_to_pdf(file: UploadFile = File(...)):
    """
    Convert PPTX to PDF using LibreOffice headless mode

    This preserves formatting and allows text selection in the browser.
    Requires LibreOffice to be installed on the system.

    Returns:
        - PDF file as bytes
        - Or error if LibreOffice is not available
    """
    logger.info(f"Converting PPTX to PDF: {file.filename}")

    import subprocess
    import shutil

    # Check if LibreOffice is available
    libreoffice_paths = []
    if os.name == 'nt':  # Windows
        libreoffice_paths = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
            r"C:\Program Files\LibreOffice\program\soffice.com",
        ]
    else:  # Linux/Mac
        libreoffice_paths = [
            "/usr/bin/soffice",
            "/usr/bin/libreoffice",
            "/Applications/LibreOffice.app/Contents/MacOS/soffice"
        ]

    # Find LibreOffice executable
    soffice_path = None
    for path in libreoffice_paths:
        if os.path.exists(path):
            soffice_path = path
            break
    else:
        # Try finding in PATH
        soffice_path = shutil.which('soffice') or shutil.which('libreoffice')

    if not soffice_path:
        logger.warning("LibreOffice not found. PPTX to PDF conversion unavailable.")
        raise HTTPException(
            status_code=503,
            detail="LibreOffice is not installed. Cannot convert PPTX to PDF. Please install LibreOffice or use the text extraction endpoint."
        )

    try:
        # Create temporary directory for conversion
        with tempfile.TemporaryDirectory() as temp_dir:
            # Save uploaded PPTX to temp file
            pptx_path = os.path.join(temp_dir, file.filename)
            with open(pptx_path, 'wb') as f:
                content = await file.read()
                f.write(content)

            logger.info(f"Saved PPTX to: {pptx_path}")

            # Convert PPTX to PDF using LibreOffice
            # --headless: run without GUI
            # --convert-to pdf: convert to PDF format
            # --outdir: output directory
            cmd = [
                soffice_path,
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', temp_dir,
                pptx_path
            ]

            logger.info(f"Running LibreOffice conversion: {' '.join(cmd)}")

            # Run conversion with timeout
            process = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                timeout=60  # 60 second timeout
            )

            if process.returncode != 0:
                logger.error(f"LibreOffice conversion failed: {process.stderr}")
                raise HTTPException(
                    status_code=500,
                    detail=f"PDF conversion failed: {process.stderr}"
                )

            # Find generated PDF file
            pdf_filename = Path(file.filename).stem + '.pdf'
            pdf_path = os.path.join(temp_dir, pdf_filename)

            if not os.path.exists(pdf_path):
                logger.error(f"PDF file not generated: {pdf_path}")
                raise HTTPException(
                    status_code=500,
                    detail="PDF conversion failed - output file not found"
                )

            # Read PDF file
            with open(pdf_path, 'rb') as pdf_file:
                pdf_bytes = pdf_file.read()

            logger.info(f"Successfully converted PPTX to PDF: {len(pdf_bytes)} bytes")

            # Return PDF as response
            from fastapi.responses import Response
            return Response(
                content=pdf_bytes,
                media_type='application/pdf',
                headers={
                    'Content-Disposition': f'attachment; filename="{pdf_filename}"'
                }
            )

    except subprocess.TimeoutExpired:
        logger.error("LibreOffice conversion timeout")
        raise HTTPException(
            status_code=500,
            detail="PDF conversion timeout - file may be too large or complex"
        )
    except Exception as e:
        logger.error(f"Error converting PPTX to PDF: {str(e)}\n{traceback.format_exc()}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to convert PPTX to PDF: {str(e)}"
        )


@app.post("/extract-images")
async def extract_images(file: UploadFile = File(...)) -> JSONResponse:
    """
    Extract images from PDF files and save them to temporary files

    Returns:
        - images: List of extracted images with paths and page numbers
        - total_images: Number of images found
        - success: Processing status
        - filename: Original filename
    """
    logger.info(f"Extracting images from PDF: {file.filename}")

    try:
        # Read file content
        content = await file.read()

        # Create temporary directory for images
        temp_dir = tempfile.mkdtemp(prefix='pdf_images_')
        logger.info(f"Created temp directory: {temp_dir}")

        extracted_images = []

        # Extract images using pdfplumber
        with pdfplumber.open(io.BytesIO(content)) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                # Get images from page
                if hasattr(page, 'images') and page.images:
                    for img_index, img in enumerate(page.images, 1):
                        try:
                            # Extract image bbox
                            x0, top, x1, bottom = img['x0'], img['top'], img['x1'], img['bottom']

                            # Crop the image from page
                            cropped_img = page.within_bbox((x0, top, x1, bottom))

                            # Convert to PIL Image
                            img_obj = cropped_img.to_image(resolution=150)
                            pil_img = img_obj.original

                            # Save image to temp file
                            img_filename = f"page_{page_num}_img_{img_index}.png"
                            img_path = os.path.join(temp_dir, img_filename)
                            pil_img.save(img_path, 'PNG')

                            extracted_images.append({
                                "imagePath": img_path,
                                "pageNumber": page_num,
                                "imageIndex": img_index
                            })

                            logger.debug(f"Extracted image {img_index} from page {page_num}")

                        except Exception as img_error:
                            logger.warning(f"Failed to extract image {img_index} from page {page_num}: {str(img_error)}")
                            continue

        if not extracted_images:
            logger.info(f"No images found in PDF: {file.filename}")
            return JSONResponse(
                status_code=200,
                content={
                    "images": [],
                    "total_images": 0,
                    "success": True,
                    "message": "No images found in document",
                    "filename": file.filename
                }
            )

        logger.info(f"Successfully extracted {len(extracted_images)} images from {file.filename}")

        return JSONResponse(
            status_code=200,
            content={
                "images": extracted_images,
                "total_images": len(extracted_images),
                "success": True,
                "filename": file.filename,
                "temp_dir": temp_dir
            }
        )

    except Exception as e:
        logger.error(f"Error extracting images from {file.filename}: {str(e)}\n{traceback.format_exc()}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to extract images: {str(e)}"
        )


@app.post("/extract/auto")
async def extract_auto(file: UploadFile = File(...)) -> JSONResponse:
    """
    Auto-detect file type and extract text accordingly

    Routes to appropriate extraction endpoint based on file extension
    """
    logger.info(f"Auto-detecting and processing: {file.filename}")

    # Get file extension
    file_ext = Path(file.filename).suffix.lower()

    # Route to appropriate handler
    if file_ext == '.pdf':
        return await extract_pdf(file)
    elif file_ext == '.docx':
        return await extract_docx(file)
    elif file_ext == '.pptx':
        return await extract_pptx(file)
    elif file_ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.tif']:
        return await extract_image_ocr(file)
    else:
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported file type: {file_ext}. Supported: .pdf, .docx, .pptx, .jpg, .png, .gif, .bmp, .tiff"
        )


@app.post("/ocr")
async def comprehensive_ocr(file: UploadFile = File(...)) -> JSONResponse:
    """
    Comprehensive OCR endpoint for managed RAG
    Handles both images and scanned PDFs
    Returns structured text with page markers

    Supports:
        - Images: JPG, JPEG, PNG, GIF, BMP, TIFF
        - Scanned PDFs: Multi-page PDFs converted to images and OCR'd

    Returns:
        - text: Extracted text with page markers
        - pageCount: Number of pages processed
        - language: Detected language (default: 'en')
        - confidence: Average OCR confidence
        - success: Processing status
    """
    logger.info(f"Processing file with OCR: {file.filename}")

    try:
        # Check if Tesseract is available
        if not _check_tesseract():
            raise HTTPException(
                status_code=503,
                detail="Tesseract OCR is not available. Please install Tesseract-OCR."
            )

        # Read file content
        content = await file.read()
        file_ext = Path(file.filename).suffix.lower()

        extracted_texts = []
        total_confidence = []

        # Handle images
        if file_ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.tif']:
            logger.info(f"Processing image file: {file.filename}")

            # Open image with PIL
            image = Image.open(io.BytesIO(content))

            # Convert to RGB if necessary
            if image.mode not in ('RGB', 'L'):
                image = image.convert('RGB')

            # Resize if too large
            max_dimension = 3000
            if max(image.size) > max_dimension:
                ratio = max_dimension / max(image.size)
                new_size = tuple(int(dim * ratio) for dim in image.size)
                image = image.resize(new_size, Image.Resampling.LANCZOS)
                logger.info(f"Resized image to {new_size}")

            # Perform OCR
            text = pytesseract.image_to_string(image, lang='eng')
            extracted_texts.append(text)

            # Get confidence
            try:
                ocr_data = pytesseract.image_to_data(image, output_type=pytesseract.Output.DICT)
                confidences = [conf for conf in ocr_data['conf'] if conf != -1]
                if confidences:
                    total_confidence.append(sum(confidences) / len(confidences))
            except Exception:
                pass

            page_count = 1

        # Handle PDFs
        elif file_ext == '.pdf':
            logger.info(f"Processing PDF file: {file.filename}")

            # Try to extract text first (to detect if it's already text-based)
            try:
                with pdfplumber.open(io.BytesIO(content)) as pdf:
                    total_text = ""
                    for page in pdf.pages:
                        text = page.extract_text() or ""
                        total_text += text

                    # If we got substantial text, it's not scanned
                    avg_chars_per_page = len(total_text) / len(pdf.pages)
                    if avg_chars_per_page > 100:
                        logger.info(f"PDF appears to be text-based ({avg_chars_per_page:.0f} chars/page). Not using OCR.")
                        return JSONResponse(
                            status_code=200,
                            content={
                                "success": False,
                                "needsOCR": False,
                                "message": "PDF has extractable text, OCR not needed"
                            }
                        )
            except Exception:
                pass

            # If we're here, it's a scanned PDF - convert pages to images and OCR
            try:
                from pdf2image import convert_from_bytes

                logger.info("Converting PDF pages to images for OCR...")

                # Convert PDF to images
                images = convert_from_bytes(content, dpi=300)
                page_count = len(images)

                logger.info(f"Processing {page_count} pages with OCR...")

                for page_num, image in enumerate(images, start=1):
                    # Convert to RGB if necessary
                    if image.mode not in ('RGB', 'L'):
                        image = image.convert('RGB')

                    # Perform OCR
                    text = pytesseract.image_to_string(image, lang='eng')

                    # Add page marker
                    page_marker = f"\n\n--- Page {page_num} ---\n\n"
                    extracted_texts.append(page_marker + text)

                    # Get confidence
                    try:
                        ocr_data = pytesseract.image_to_data(image, output_type=pytesseract.Output.DICT)
                        confidences = [conf for conf in ocr_data['conf'] if conf != -1]
                        if confidences:
                            total_confidence.append(sum(confidences) / len(confidences))
                    except Exception:
                        pass

                    logger.info(f"Processed page {page_num}/{page_count}")

            except ImportError:
                raise HTTPException(
                    status_code=503,
                    detail="pdf2image library not available. Install with: pip install pdf2image"
                )

        else:
            raise HTTPException(
                status_code=400,
                detail=f"Unsupported file type for OCR: {file_ext}"
            )

        # Combine all text
        combined_text = "".join(extracted_texts)

        if not combined_text.strip():
            logger.warning(f"No text extracted from {file.filename}")
            return JSONResponse(
                status_code=200,
                content={
                    "text": "",
                    "success": True,
                    "pageCount": page_count,
                    "warning": "No text detected",
                    "filename": file.filename
                }
            )

        # Calculate average confidence
        avg_confidence = sum(total_confidence) / len(total_confidence) if total_confidence else None

        logger.info(f"Successfully extracted {len(combined_text)} characters via OCR from {page_count} page(s)")

        result = {
            "text": combined_text,
            "success": True,
            "pageCount": page_count,
            "language": "en",  # Could be enhanced with language detection
            "filename": file.filename,
            "charCount": len(combined_text),
            "method": "ocr"
        }

        if avg_confidence is not None:
            result["confidence"] = round(avg_confidence, 2)

        return JSONResponse(status_code=200, content=result)

    except Exception as e:
        logger.error(f"Error in OCR processing {file.filename}: {str(e)}\n{traceback.format_exc()}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to process file with OCR: {str(e)}"
        )


if __name__ == "__main__":
    import uvicorn

    logger.info("Starting Document Processing Service...")
    logger.info("Service will be available at: http://localhost:8000")
    logger.info("API docs available at: http://localhost:8000/docs")

    uvicorn.run(
        app,
        host="0.0.0.0",
        port=8000,
        log_level="info"
    )